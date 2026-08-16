#!/usr/bin/env python3
"""
Phase -1 visual lab for the PowerPoint Python-highlighting add-in.

Generates a real .pptx with one slide per (font, size, line-spacing) variant so
the look can be judged in PowerPoint itself rather than approximated in HTML.

Whichever slide wins becomes the spec that the VBA renderer must reproduce.

Usage:  python3 tools/lab.py [-o OUT.pptx]
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt
from pygments import lex
from pygments.lexers import PythonLexer
from pygments.token import Comment, Keyword, Name, Number, Operator, String, Token

REPO = Path(__file__).resolve().parent.parent
# Samples are filed per language, so a second language is a sibling directory
# rather than a naming convention inside a shared one.
SAMPLES = REPO / "tests" / "samples" / "python"
SNIPPET = SAMPLES / "lab_snippet.py"

# --------------------------------------------------------------------------
# Palette - VS Code "Dark Modern" / Dark+ token colors
# --------------------------------------------------------------------------

PALETTE = {
    "bg":       "1F1F1F",   # editor background
    "default":  "D4D4D4",   # operators, punctuation
    "comment":  "6A9955",
    "string":   "CE9178",
    "number":   "B5CEA8",
    "kw_ctrl":  "C586C0",   # if / for / return / import
    "kw_decl":  "569CD6",   # def / class / None / self / and / not
    "func":     "DCDCAA",   # definitions, call sites, builtins, decorators
    "cls":      "4EC9B0",   # class names, types
    "var":      "9CDCFE",   # variables, parameters, attributes
    # Bracket pairs by nesting depth, cycling every three levels. On by default
    # in VS Code, and on a line like print(id(x)) it is most of what you see.
    "br1":      "FFD700",   # gold
    "br2":      "DA70D6",   # orchid
    "br3":      "179FFF",   # blue
    "gutter":   "858585",
    "label":    "808080",
}

# Softer variant, in case Dark+ is too saturated on a classroom projector.
PALETTE_MUTED = dict(PALETTE, **{
    "string":  "C99076",
    "kw_ctrl": "BC8FB8",
    "kw_decl": "6B9FD0",
    "func":    "CFCF9E",
    "var":     "9BC4DC",
})

CONTROL_KW = {
    "if", "elif", "else", "for", "while", "return", "break", "continue",
    "try", "except", "finally", "raise", "with", "as", "yield", "pass",
    "assert", "del", "import", "from", "global", "nonlocal", "async",
    "await", "match", "case",
}

# --------------------------------------------------------------------------
# Fonts - advance and x-height ratios measured from the actual .ttf files
# --------------------------------------------------------------------------

FONTS = {
    # name shown in PowerPoint : (advance/em, x-height/em, note)
    "Consolas":          (0.550, 0.490, "VS Code default on Windows"),
    "Cascadia Mono":     (0.586, 0.518, "Windows 11 / Terminal"),
    "JetBrains Mono NL": (0.600, 0.550, "PyCharm default"),
}

SIZES = [18, 20, 22, 24]        # the classroom-projection band
SPACING_RATIOS = [1.20, 1.35]   # 1.35 is VS Code's own default line height

# --------------------------------------------------------------------------
# LOCKED SPEC (chosen from the lab, slide 5)
#
# Font size is the ONLY free parameter. Everything else derives from it, so a
# block can be resized on any slide without any of its proportions breaking.
# The VBA renderer must use exactly these formulas.
# --------------------------------------------------------------------------

BASE_SIZE = 22.0            # default pt size for a new code block
SIZE_LADDER = [10, 12, 14, 16, 18, 20, 22, 24, 28, 32]   # A+ / A- steps
MIN_TEACHING_SIZE = 16      # below this, warn - it will not read from the back

LINE_RATIO = 1.20           # line spacing = 1.20 x size, set in EXACT points
PAD_RATIO = 0.64            # internal padding  = 0.64 x size  (14.1pt at 22pt)
GAP_RATIO = 0.45            # gutter-to-code gap = 0.45 x size
RADIUS_RATIO = 0.36         # corner radius     = 0.36 x size  (8pt at 22pt)
ADVANCE = 0.550             # Consolas, measured from consola.ttf


def derive(size: float, n_lines: int):
    """All block geometry as a function of font size. The whole spec."""
    digits = max(2, len(str(max(1, n_lines))))
    return {
        "size":    size,
        "line":    round(size * LINE_RATIO, 1),
        "pad":     round(size * PAD_RATIO, 1),
        "gutter":  round(digits * ADVANCE * size + size * GAP_RATIO, 1),
        "radius":  round(size * RADIUS_RATIO, 1),
        "charw":   size * ADVANCE,
    }


def fit_size(n_lines: int, max_chars: int, avail_w_pt: float, avail_h_pt: float):
    """Largest ladder size at which the snippet fits the available area.

    Drives the 'Fit' ribbon button, and also caps A+ so growing a block can
    never push code off the edge of the slide.
    """
    best = SIZE_LADDER[0]
    for s in SIZE_LADDER:
        d = derive(s, n_lines)
        w = max_chars * d["charw"] + d["gutter"] + 2 * d["pad"]
        h = n_lines * d["line"] + 2 * d["pad"]
        if w <= avail_w_pt and h <= avail_h_pt:
            best = s
    return best


def split_advice(n_lines: int, max_chars: int, avail_w_pt: float, avail_h_pt: float):
    """How many slides this snippet needs to stay readable from the back.

    Returns (fitting_size, slides_needed_at_min_teaching_size).
    """
    fitted = fit_size(n_lines, max_chars, avail_w_pt, avail_h_pt)
    if fitted >= MIN_TEACHING_SIZE:
        return fitted, 1
    d = derive(MIN_TEACHING_SIZE, n_lines)
    per_slide = max(1, int((avail_h_pt - 2 * d["pad"]) / d["line"]))
    return fitted, -(-n_lines // per_slide)   # ceiling division

# --------------------------------------------------------------------------
# Slide geometry (16:9, 13.333in x 7.5in)
# --------------------------------------------------------------------------

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
BLOCK_L, BLOCK_T = Inches(0.55), Inches(1.05)
BLOCK_W, BLOCK_H = Inches(12.23), Inches(5.95)
PAD = Pt(14)            # internal padding inside the dark block
GUTTER_W = Inches(0.55)


BRACKET_KEYS = ("br1", "br2", "br3")
OPEN_BRACKETS = "([{"
CLOSE_BRACKETS = ")]}"


def apply_bracket_depth(runs: list[tuple[str, str]]) -> list[tuple[str, str]]:
    """Split brackets out of punctuation runs and colour them by nesting depth.

    Only runs already classified as `default` are examined, so brackets inside
    strings and comments are left alone - they are not structure.
    """
    out: list[tuple[str, str]] = []
    depth = 0
    for key, value in runs:
        if key != "default":
            out.append((key, value))
            continue
        buf = ""
        for ch in value:
            if ch in OPEN_BRACKETS or ch in CLOSE_BRACKETS:
                if buf:
                    out.append(("default", buf))
                    buf = ""
                if ch in OPEN_BRACKETS:
                    out.append((BRACKET_KEYS[depth % 3], ch))
                    depth += 1
                else:
                    # Decrement first so a closer matches its opener's colour.
                    depth = max(0, depth - 1)
                    out.append((BRACKET_KEYS[depth % 3], ch))
            else:
                buf += ch
        if buf:
            out.append(("default", buf))
    return out


def classify(tokens: list[tuple]) -> list[tuple[str, str]]:
    """Map pygments tokens onto our palette keys, VS Code style."""
    out: list[tuple[str, str]] = []
    # A decorator is only a decorator at the start of a line. Whitespace keeps
    # the flag, so an indented @property still counts.
    at_line_start = True
    for i, (ttype, value) in enumerate(tokens):
        # look ahead past whitespace for a '(' to detect call sites
        nxt = ""
        for j in range(i + 1, len(tokens)):
            if tokens[j][1].strip():
                nxt = tokens[j][1].lstrip()
                break

        if ttype in Comment:
            key = "comment"
        elif ttype in String:
            key = "string"
        elif ttype in Number:
            key = "number"
        elif ttype in Name and value in ("j", "J") and out and out[-1][0] == "number":
            # pygments splits a complex literal: 3j lexes as Number "3" plus
            # Name "j". VS Code colours the whole literal as a number, so glue
            # the suffix back on.
            key = "number"
        elif ttype in Name.Decorator and not at_line_start:
            # pygments matches @ as a decorator anywhere, so the matrix-multiply
            # operator in "a @ b" comes out coloured like @decorator. Only a @
            # that starts a line is a decorator. Split the token: the operator
            # is punctuation, the name after it is an ordinary name.
            out.append(("default", value[:1]))
            rest = value[1:]
            if rest:
                out.append(("func" if nxt.startswith("(") else "var", rest))
            at_line_start = False
            continue
        elif ttype in Keyword:
            key = "kw_ctrl" if value in CONTROL_KW else "kw_decl"
        elif ttype in Name.Class:
            key = "cls"
        elif ttype in Name.Function or ttype in Name.Decorator:
            key = "func"
        elif ttype in Name.Builtin:
            key = "func" if nxt.startswith("(") else "kw_decl"
        elif ttype in Name:
            # a bare name followed by '(' is a call, otherwise a variable
            key = "func" if nxt.startswith("(") else "var"
        elif ttype in Operator.Word:
            # and / or / not / in / is - VS Code renders these blue, but
            # pygments files them under Operator rather than Keyword.
            key = "kw_decl"
        elif ttype in Operator or ttype in Token.Punctuation:
            key = "default"
        else:
            key = "default"
        out.append((key, value))
        if value.strip():
            at_line_start = value.endswith("\n")
        elif "\n" in value:
            at_line_start = True
    return apply_bracket_depth(out)


def to_lines(classified: list[tuple[str, str]]) -> list[list[tuple[str, str]]]:
    """Split (key, text) runs into one list of runs per source line."""
    lines: list[list[tuple[str, str]]] = [[]]
    for key, value in classified:
        parts = value.split("\n")
        for n, part in enumerate(parts):
            if n:
                lines.append([])
            if part:
                lines[-1].append((key, part))
    while lines and not lines[-1]:
        lines.pop()
    return lines


def add_code_block(slide, lines, font, size, spacing_pt, palette, gutter):
    """Draw the dark block, with the line-number gutter inside it.

    The block hugs its content vertically rather than filling a fixed height,
    and the gutter overlays the block's left edge so the numbers sit on the
    dark background exactly as they do in VS Code.
    """
    d = derive(size, len(lines))
    pad, gutter_w = Pt(d["pad"]), Pt(d["gutter"])

    height = Pt(len(lines) * spacing_pt + 2 * d["pad"])
    top = Emu(int(BLOCK_T + (BLOCK_H - height) / 2))   # centre in the content area
    left, width = BLOCK_L, BLOCK_W

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # adjustment is a fraction of the shorter side, so hold the radius constant
    box.adjustments[0] = min(0.5, d["radius"] / max(1.0, Emu(height).inches * 72))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string(palette["bg"])
    box.line.fill.background()
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Emu(pad + gutter_w) if gutter else pad
    tf.margin_right = pad
    tf.margin_top = tf.margin_bottom = pad

    for n, runs in enumerate(lines):
        para = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
        # Autoshapes default to centered text. Code must be explicitly left
        # aligned or every line floats. The VBA renderer must do this too.
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = Pt(spacing_pt)
        para.space_before = Pt(0)
        para.space_after = Pt(0)
        for key, text in runs or [("default", "")]:
            run = para.add_run()
            run.text = text.replace("\t", "    ")
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = RGBColor.from_string(palette[key])

    if gutter:
        g = slide.shapes.add_textbox(left, top, Emu(pad + gutter_w), height)
        gtf = g.text_frame
        gtf.word_wrap = False
        gtf.auto_size = MSO_AUTO_SIZE.NONE
        gtf.margin_left = Pt(0)
        gtf.margin_right = Pt(round(size * GAP_RATIO, 1))
        gtf.margin_top = gtf.margin_bottom = pad
        for n in range(len(lines)):
            para = gtf.paragraphs[0] if n == 0 else gtf.add_paragraph()
            para.line_spacing = Pt(spacing_pt)
            para.space_before = para.space_after = Pt(0)
            para.alignment = PP_ALIGN.RIGHT
            run = para.add_run()
            run.text = str(n + 1)
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = RGBColor.from_string(palette["gutter"])
    return box


def add_label(slide, text):
    tb = slide.shapes.add_textbox(BLOCK_L, Inches(0.28), BLOCK_W, Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Segoe UI"
    r.font.size = Pt(13)
    r.font.color.rgb = RGBColor.from_string(PALETTE["label"])


def build(out_path: Path):
    code = SNIPPET.read_text()
    lines = to_lines(classify(list(lex(code, PythonLexer()))))

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    inner_w_pt = Emu(BLOCK_W - GUTTER_W).inches * 72 - 2 * PAD.pt
    inner_h_pt = Emu(BLOCK_H).inches * 72 - 2 * PAD.pt

    variants = []
    for font, (adv, xh, note) in FONTS.items():
        for size in SIZES:
            for ratio in SPACING_RATIOS:
                variants.append((font, adv, xh, note, size, ratio))

    for font, adv, xh, note, size, ratio in variants:
        spacing_pt = round(size * ratio, 1)
        cols = int(inner_w_pt / (size * adv))
        rows = int(inner_h_pt / spacing_pt)
        xheight_pt = round(size * xh, 1)

        slide = prs.slides.add_slide(blank)
        add_label(slide, (
            f"{font}  ·  {size}pt  ·  line {spacing_pt}pt ({ratio:g}×)  "
            f"·  fits {rows} lines × {cols} chars  "
            f"·  x-height {xheight_pt}pt  ·  {note}"
        ))
        add_code_block(slide, lines, font, size, spacing_pt, PALETTE, gutter=True)

    # One muted-palette comparison at a mid setting, and one gutter-off.
    for pal, gutter, tag in ((PALETTE_MUTED, True, "muted palette"),
                             (PALETTE, False, "no line numbers")):
        slide = prs.slides.add_slide(blank)
        add_label(slide, f"Consolas  ·  20pt  ·  line 24pt (1.2×)  ·  {tag}")
        add_code_block(slide, lines, "Consolas", 20, 24.0, pal, gutter=gutter)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote {out_path}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


def build_ladder(out_path: Path):
    """Verify the locked spec holds across the whole size ladder."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    blank = prs.slide_layouts[6]

    avail_w_pt = Emu(BLOCK_W).inches * 72
    avail_h_pt = Emu(BLOCK_H).inches * 72

    short = to_lines(classify(list(lex(SNIPPET.read_text(), PythonLexer()))))
    long_src = (SAMPLES / "long_snippet.py").read_text()
    long_lines = to_lines(classify(list(lex(long_src, PythonLexer()))))

    for size in SIZE_LADDER:
        d = derive(size, len(short))
        slide = prs.slides.add_slide(blank)
        tag = "  <- DEFAULT (slide 5)" if size == BASE_SIZE else ""
        add_label(slide, (
            f"Consolas {size}pt  ·  line {d['line']}pt  ·  pad {d['pad']}pt  "
            f"·  gutter {d['gutter']}pt  ·  radius {d['radius']}pt{tag}"
        ))
        add_code_block(slide, short, "Consolas", size, d["line"], PALETTE, gutter=True)

    # The 'Fit' button: largest ladder size at which a long snippet still fits.
    n, cols = len(long_lines), max(len(_plain(r)) for r in long_lines)
    best = fit_size(n, cols, avail_w_pt, avail_h_pt)
    d = derive(best, n)
    slide = prs.slides.add_slide(blank)
    add_label(slide, (
        f"FIT: {n} lines × {cols} chars → auto-chose Consolas {best}pt  "
        f"·  line {d['line']}pt  ·  gutter {d['gutter']}pt"
    ))
    add_code_block(slide, long_lines, "Consolas", best, d["line"], PALETTE, gutter=True)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"wrote {out_path}  (fit chose {best}pt for {n} lines × {cols} chars)")


def _plain(runs):
    return "".join(t for _, t in runs)


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("-o", "--out")
    ap.add_argument("--ladder", action="store_true",
                    help="render the locked spec across the size ladder")
    a = ap.parse_args()
    default = "ladder.pptx" if a.ladder else "lab.pptx"
    out = Path(a.out) if a.out else REPO / "dist" / default
    (build_ladder if a.ladder else build)(out)
